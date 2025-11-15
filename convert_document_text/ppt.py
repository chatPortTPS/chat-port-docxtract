from pathlib import Path
from typing import Dict, Any, List
from pptx import Presentation
import unicodedata


class Ppt:
    """
    Clase para extraer texto de presentaciones PowerPoint (.pptx)
    """

    def __init__(self):
        print("Inicializando extractor PowerPoint")

    def _normalize_text(self, text):
        """Normaliza el texto eliminando acentos y convirtiendo a minúsculas."""
        text = text.lower()
        text = ''.join(
            c for c in unicodedata.normalize('NFD', text)
            if unicodedata.category(c) != 'Mn'
        )
        return text

    def _extraer_metadatos_ppt(self, prs: Presentation) -> Dict[str, Any]:
        """Extrae metadatos de la presentación."""
        core_props = prs.core_properties
        return {
            'titulo': core_props.title or '',
            'autor': core_props.author or '',
            'fecha_creacion': str(core_props.created) if core_props.created else '',
            'fecha_modificacion': str(core_props.modified) if core_props.modified else '',
            'diapositivas': len(prs.slides)
        }

    def _extraer_texto_shape(self, shape) -> str:
        """Extrae texto de una forma (shape) en la diapositiva."""
        texto = ""
        if hasattr(shape, "text"):
            texto = shape.text.strip()
        return texto

    def _extraer_tablas(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Extrae tablas de las diapositivas."""
        tablas_procesadas = []
        indice_tabla = 0

        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    try:
                        tabla = shape.table
                        datos = []

                        for fila in tabla.rows:
                            fila_datos = [celda.text.strip() for celda in fila.cells]
                            if any(fila_datos):  # Solo agregar si hay contenido
                                datos.append(fila_datos)

                        if datos:
                            indice_tabla += 1
                            tablas_procesadas.append({
                                'indice': indice_tabla,
                                'diapositiva': slide_num + 1,
                                'datos': datos,
                                'filas': len(datos),
                                'columnas': len(datos[0]) if datos else 0
                            })
                    except Exception as e:
                        print(f"Error procesando tabla en diapositiva {slide_num + 1}: {e}")
                        continue

        return tablas_procesadas

    def procesar_ppt(self, ruta_archivo: str) -> Dict[str, Any]:
        """
        Procesa un archivo PowerPoint extrayendo texto, tablas y metadatos.

        Args:
            ruta_archivo: Ruta al archivo PowerPoint

        Returns:
            Diccionario con contenido estructurado de la presentación
        """
        resultado = {
            'tipo_archivo': 'powerpoint',
            'archivo': Path(ruta_archivo).name,
            'metadatos': {},
            'contenido': {
                'texto': '',
                'tablas': []
            }
        }

        # Abrir la presentación
        prs = Presentation(ruta_archivo)

        try:
            # Extraer metadatos
            resultado['metadatos'] = self._extraer_metadatos_ppt(prs)

            # Extraer texto de todas las diapositivas
            texto_completo = []

            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    texto = self._extraer_texto_shape(shape)
                    if texto:
                        texto_completo.append(texto)

            resultado['contenido']['texto'] = self._normalize_text(' '.join(texto_completo))

            # Extraer tablas
            resultado['contenido']['tablas'] = self._extraer_tablas(prs)

        except Exception as e:
            print(f"Error procesando presentación PowerPoint: {e}")
            raise

        return resultado
